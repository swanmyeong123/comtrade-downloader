import streamlit as st
import pandas as pd
import requests
import time
import datetime
from io import BytesIO

# --- 1. 국가 코드 정의 ---
# EU 27개국 리스트 (브렉시트 이후)
EU27_LIST = [
    "040", "056", "100", "191", "196", "203", "208", "233", "246", "251", 
    "276", "300", "348", "372", "380", "428", "440", "442", "470", "528", 
    "616", "620", "703", "705", "724", "752", "242"
]
EU27_STR = ",".join(EU27_LIST)

# CPTPP 등 기타 그룹
CPTPP_11_STR = "036,096,124,152,392,458,484,554,604,702,704" # 영국 미포함
UK_CODE = "826"

# UN M49 국가 코드 → 영문 국가명 매핑
COUNTRY_NAMES = {
    # World / Regions
    "0": "World",
    "002": "Africa",
    "009": "Oceania",
    "019": "Americas",
    "142": "Asia",
    "150": "Europe",
    # EU 27 Countries
    "040": "Austria",
    "056": "Belgium",
    "100": "Bulgaria",
    "191": "Croatia",
    "196": "Cyprus",
    "203": "Czechia",
    "208": "Denmark",
    "233": "Estonia",
    "246": "Finland",
    "251": "France",
    "242": "Fiji",  # Note: This might be an error in EU27_LIST, as Fiji is not in EU
    "276": "Germany",
    "300": "Greece",
    "348": "Hungary",
    "372": "Ireland",
    "380": "Italy",
    "428": "Latvia",
    "440": "Lithuania",
    "442": "Luxembourg",
    "470": "Malta",
    "528": "Netherlands",
    "616": "Poland",
    "620": "Portugal",
    "703": "Slovakia",
    "705": "Slovenia",
    "724": "Spain",
    "752": "Sweden",
    # Other Major Countries
    "004": "Afghanistan",
    "008": "Albania",
    "012": "Algeria",
    "020": "Andorra",
    "024": "Angola",
    "028": "Antigua and Barbuda",
    "031": "Azerbaijan",
    "032": "Argentina",
    "036": "Australia",
    "044": "Bahamas",
    "048": "Bahrain",
    "050": "Bangladesh",
    "051": "Armenia",
    "052": "Barbados",
    "064": "Bhutan",
    "068": "Bolivia",
    "070": "Bosnia and Herzegovina",
    "072": "Botswana",
    "076": "Brazil",
    "084": "Belize",
    "090": "Solomon Islands",
    "096": "Brunei Darussalam",
    "104": "Myanmar",
    "108": "Burundi",
    "112": "Belarus",
    "116": "Cambodia",
    "120": "Cameroon",
    "124": "Canada",
    "132": "Cabo Verde",
    "140": "Central African Republic",
    "144": "Sri Lanka",
    "148": "Chad",
    "152": "Chile",
    "156": "China",
    "158": "Taiwan",
    "170": "Colombia",
    "174": "Comoros",
    "178": "Congo",
    "180": "DR Congo",
    "184": "Cook Islands",
    "188": "Costa Rica",
    "192": "Cuba",
    "204": "Benin",
    "212": "Dominica",
    "214": "Dominican Republic",
    "218": "Ecuador",
    "222": "El Salvador",
    "226": "Equatorial Guinea",
    "231": "Ethiopia",
    "232": "Eritrea",
    "234": "Faroe Islands",
    "238": "Falkland Islands",
    "242": "Fiji",
    "250": "France",
    "254": "French Guiana",
    "258": "French Polynesia",
    "262": "Djibouti",
    "266": "Gabon",
    "268": "Georgia",
    "270": "Gambia",
    "275": "Palestine",
    "288": "Ghana",
    "292": "Gibraltar",
    "296": "Kiribati",
    "304": "Greenland",
    "308": "Grenada",
    "312": "Guadeloupe",
    "316": "Guam",
    "320": "Guatemala",
    "324": "Guinea",
    "328": "Guyana",
    "332": "Haiti",
    "340": "Honduras",
    "344": "Hong Kong",
    "352": "Iceland",
    "356": "India",
    "360": "Indonesia",
    "364": "Iran",
    "368": "Iraq",
    "376": "Israel",
    "384": "Cote d'Ivoire",
    "388": "Jamaica",
    "392": "Japan",
    "398": "Kazakhstan",
    "400": "Jordan",
    "404": "Kenya",
    "408": "North Korea",
    "410": "South Korea",
    "414": "Kuwait",
    "417": "Kyrgyzstan",
    "418": "Laos",
    "422": "Lebanon",
    "426": "Lesotho",
    "430": "Liberia",
    "434": "Libya",
    "438": "Liechtenstein",
    "446": "Macao",
    "450": "Madagascar",
    "454": "Malawi",
    "458": "Malaysia",
    "462": "Maldives",
    "466": "Mali",
    "474": "Martinique",
    "478": "Mauritania",
    "480": "Mauritius",
    "484": "Mexico",
    "492": "Monaco",
    "496": "Mongolia",
    "498": "Moldova",
    "499": "Montenegro",
    "500": "Montserrat",
    "504": "Morocco",
    "508": "Mozambique",
    "512": "Oman",
    "516": "Namibia",
    "520": "Nauru",
    "524": "Nepal",
    "530": "Netherlands Antilles",
    "531": "Curacao",
    "533": "Aruba",
    "534": "Sint Maarten",
    "540": "New Caledonia",
    "548": "Vanuatu",
    "554": "New Zealand",
    "558": "Nicaragua",
    "562": "Niger",
    "566": "Nigeria",
    "570": "Niue",
    "574": "Norfolk Island",
    "578": "Norway",
    "580": "Northern Mariana Islands",
    "583": "Micronesia",
    "584": "Marshall Islands",
    "585": "Palau",
    "586": "Pakistan",
    "591": "Panama",
    "598": "Papua New Guinea",
    "600": "Paraguay",
    "604": "Peru",
    "608": "Philippines",
    "612": "Pitcairn",
    "630": "Puerto Rico",
    "634": "Qatar",
    "638": "Reunion",
    "642": "Romania",
    "643": "Russia",
    "646": "Rwanda",
    "654": "Saint Helena",
    "659": "Saint Kitts and Nevis",
    "660": "Anguilla",
    "662": "Saint Lucia",
    "666": "Saint Pierre and Miquelon",
    "670": "Saint Vincent and the Grenadines",
    "674": "San Marino",
    "678": "Sao Tome and Principe",
    "682": "Saudi Arabia",
    "686": "Senegal",
    "688": "Serbia",
    "690": "Seychelles",
    "694": "Sierra Leone",
    "702": "Singapore",
    "704": "Vietnam",
    "706": "Somalia",
    "710": "South Africa",
    "716": "Zimbabwe",
    "720": "Yemen",
    "728": "South Sudan",
    "729": "Sudan",
    "732": "Western Sahara",
    "740": "Suriname",
    "748": "Eswatini",
    "756": "Switzerland",
    "760": "Syria",
    "762": "Tajikistan",
    "764": "Thailand",
    "768": "Togo",
    "772": "Tokelau",
    "776": "Tonga",
    "780": "Trinidad and Tobago",
    "784": "United Arab Emirates",
    "788": "Tunisia",
    "792": "Turkey",
    "795": "Turkmenistan",
    "796": "Turks and Caicos Islands",
    "798": "Tuvalu",
    "800": "Uganda",
    "804": "Ukraine",
    "807": "North Macedonia",
    "818": "Egypt",
    "826": "United Kingdom",
    "831": "Guernsey",
    "832": "Jersey",
    "833": "Isle of Man",
    "834": "Tanzania",
    "840": "United States",
    "842": "United States",
    "850": "US Virgin Islands",
    "854": "Burkina Faso",
    "858": "Uruguay",
    "860": "Uzbekistan",
    "862": "Venezuela",
    "876": "Wallis and Futuna",
    "882": "Samoa",
    "887": "Yemen",
    "894": "Zambia",
    # Special codes
    "EXTRA_EU": "EU27 Extra (Calculated)",
    "all": "All Countries"
}

# [보고 국가(Reporter) 그룹]
REPORTER_GROUPS = {
    "직접 입력 (Custom)": "",

    "폴란드 (Poland)": "616",
    "독일 (Germany)": "276",
    "스페인 (Spain)": "724",
    "벨기에 (Belgium)": "056",
    "스웨덴 (Sweden)": "752",
    "한국 (Korea)": "410",
    "EU 27 전체 (All EU Members)": EU27_STR,
    "중국 (China)": "156",
    "미국 (USA)": "842",
}

# [상대국(Partner) 그룹]
# 대륙별 국가 코드 (UN M49 기준)
CONTINENT_EUROPE = ["040","056","100","191","196","203","208","233","246","251","276","300","348","372","380","428","440","442","470","528","616","620","642","703","705","724","752","826","008","020","070","112","268","292","352","438","492","498","499","578","643","674","688","756","804","807"]
CONTINENT_AFRICA = ["012","024","072","108","120","132","140","148","174","178","180","204","226","231","232","262","266","270","288","324","384","404","426","430","450","454","466","478","480","504","508","516","562","566","646","686","690","694","706","710","716","728","729","732","748","768","788","800","834","854","894"]
CONTINENT_MIDDLE_EAST = ["048","364","368","376","400","414","422","512","634","682","760","784","887","275"]
CONTINENT_EAST_ASIA = ["156","392","410","158","344","446","408","496"]
CONTINENT_SOUTHEAST_ASIA = ["096","104","116","360","418","458","608","702","704","764"]
CONTINENT_NORTH_AMERICA = ["124","840","842"]
CONTINENT_CENTRAL_SOUTH_AMERICA = ["032","044","052","068","076","084","152","170","188","192","212","214","218","222","308","320","328","332","340","388","484","558","591","600","604","659","662","670","740","780","858","862"]
CONTINENT_OCEANIA = ["036","090","242","296","520","540","548","554","583","584","585","598","776","798","882"]

# 모든 대륙 국가를 합친 리스트 (대륙별 그룹화용)
ALL_CONTINENT_CODES = (CONTINENT_EUROPE + CONTINENT_AFRICA + CONTINENT_MIDDLE_EAST + 
                       CONTINENT_EAST_ASIA + CONTINENT_SOUTHEAST_ASIA + CONTINENT_NORTH_AMERICA + 
                       CONTINENT_CENTRAL_SOUTH_AMERICA + CONTINENT_OCEANIA)

PARTNER_GROUPS = {
    "직접 입력 (Custom)": "",
    "--- 대륙별 선택 ---": "SEPARATOR",
    "🌐 대륙별 통합 (All Continents)": "ALL_CONTINENTS",
    "🌍 유럽 (Europe)": ",".join(CONTINENT_EUROPE),
    "🌍 아프리카 (Africa)": ",".join(CONTINENT_AFRICA),
    "🌍 중동 (Middle East)": ",".join(CONTINENT_MIDDLE_EAST),
    "🌍 동아시아 (East Asia)": ",".join(CONTINENT_EAST_ASIA),
    "🌍 동남아시아 (Southeast Asia)": ",".join(CONTINENT_SOUTHEAST_ASIA),
    "🌍 북미 (North America)": ",".join(CONTINENT_NORTH_AMERICA),
    "🌍 중남미 (Central/South America)": ",".join(CONTINENT_CENTRAL_SOUTH_AMERICA),
    "🌍 오세아니아 (Oceania)": ",".join(CONTINENT_OCEANIA),
    "--- 기존 선택 ---": "SEPARATOR",
    "★ EU 27 역외 (Extra-EU) [World - EU27]": "EXTRA_EU_CALC", 
    "전 세계 합계 (World Total)": "0",
    "EU 27 (역내 교역)": EU27_STR,
    "CPTPP (11개국 - 영국 미포함)": CPTPP_11_STR,
    "CPTPP (12개국 - 영국 포함)": CPTPP_11_STR + "," + UK_CODE,
    "미국 (USA)": "842",
    "중국 (China)": "156"
}

# 대륙 이름 매핑 (코드 → 대륙명)
def get_continent_name(country_code):
    """국가 코드로부터 대륙명 반환"""
    code = str(country_code).zfill(3)
    if code in CONTINENT_EUROPE: return "Europe"
    if code in CONTINENT_AFRICA: return "Africa"
    if code in CONTINENT_MIDDLE_EAST: return "Middle East"
    if code in CONTINENT_EAST_ASIA: return "East Asia"
    if code in CONTINENT_SOUTHEAST_ASIA: return "Southeast Asia"
    if code in CONTINENT_NORTH_AMERICA: return "North America"
    if code in CONTINENT_CENTRAL_SOUTH_AMERICA: return "Central/South America"
    if code in CONTINENT_OCEANIA: return "Oceania"
    return "Others"


# UN M49 기반 국가 목록 (지역별 그룹화)
COUNTRIES_BY_REGION = {
    "Africa": {
        "Northern Africa": ["012", "818", "434", "504", "729", "788", "732"],
        "Sub-Saharan Africa": {
            "Eastern Africa": ["108", "086", "174", "262", "232", "231", "404", "175", "454", "480", "508", "638", "646", "690", "706", "728", "800", "834", "716", "894"],
            "Middle Africa": ["024", "120", "140", "148", "266", "226", "178", "180", "678"],
            "Southern Africa": ["072", "748", "426", "516", "710"],
            "Western Africa": ["204", "132", "270", "288", "324", "624", "430", "466", "478", "562", "566", "654", "686", "694", "768", "854"]
        }
    },
    "Americas": {
        "Northern America": ["060", "124", "304", "666", "840"],  
        "Latin America and the Caribbean": {
            "Caribbean": ["660", "028", "533", "044", "052", "535", "136", "192", "531", "212", "214", "308", "312", "332", "388", "474", "500", "630", "652", "659", "662", "663", "670", "780", "796", "092", "850", "534"],
            "Central America": ["084", "188", "222", "320", "340", "484", "558", "591"],
            "South America": ["032", "068", "074", "076", "170", "218", "238", "254", "328", "600", "604", "239", "740", "858", "862"]
        }
    },
    "Asia": {
        "Central Asia": ["398", "417", "762", "795", "860"],
        "Eastern Asia": ["156", "344", "392", "408", "410", "446", "496", "158"],
        "Southern Asia": ["004", "050", "064", "356", "364", "462", "524", "586", "144"],
        "South-eastern Asia": ["096", "104", "116", "360", "418", "458", "608", "702", "764", "626", "704"],
        "Western Asia": ["051", "031", "048", "196", "268", "368", "376", "400", "414", "422", "275", "512", "634", "682", "760", "784", "792", "887"]
    },
    "Europe": {
        "Eastern Europe": ["112", "100", "203", "348", "616", "498", "642", "643", "703", "804", "807"],
        "Northern Europe": ["248", "208", "233", "234", "372", "352", "428", "440", "578", "744", "752", "826", "831", "832", "833"],
        "Southern Europe": ["008", "020", "070", "191", "292", "300", "336", "380", "470", "499", "620", "674", "688", "705", "724"],
        "Western Europe": ["040", "056", "250", "276", "438", "442", "492", "528", "756"]
    },
    "Oceania": {
        "Australia and New Zealand": ["036", "162", "166", "334", "554", "574"],
        "Melanesia": ["242", "540", "598", "090", "548"],
        "Micronesia": ["316", "296", "584", "583", "520", "580", "585"],
        "Polynesia": ["016", "184", "258", "570", "612", "772", "776", "798", "882", "876"]
    }
}


API_URL = "https://comtradeapi.un.org/data/v1/get/C/A/HS"
current_year = datetime.datetime.now().year
YEAR_OPTIONS = [str(y) for y in range(current_year, 1999, -1)]

def get_comtrade_data(api_key, hs_code, single_year, reporter_code, partner_code, flow_code):
    headers = {"Ocp-Apim-Subscription-Key": api_key}
    
    # EU 역외 계산 모드일 경우: World(0)와 EU27 국가들을 모두 요청
    if partner_code == "EXTRA_EU_CALC":
        actual_partner = "0," + EU27_STR
    elif partner_code == "ALL_CONTINENTS":
        actual_partner = ",".join(ALL_CONTINENT_CODES)
    else:
        actual_partner = partner_code

    clean_reporter = reporter_code.replace(" ", "")
    clean_partner = actual_partner.replace(" ", "")
    
    params = {
        "reporterCode": clean_reporter, 
        "partnerCode": clean_partner,
        "period": single_year,
        "cmdCode": str(hs_code).strip(),
        "flowCode": flow_code,
        "motCode": "0",
        "freqCode": "A",
        "format": "json"
    }

    try:
        response = requests.get(API_URL, headers=headers, params=params, timeout=60)
        response.raise_for_status()
        
        data = response.json()
        if 'data' in data and len(data['data']) > 0:
            df = pd.DataFrame(data['data'])
            
            # EU 역외 교역 계산 (World - EU_Sum)
            if partner_code == "EXTRA_EU_CALC":
                return calculate_extra_eu(df)
            else:
                return df
        else:
            return pd.DataFrame()
            
    except Exception as e:
        print(f"Error (HS:{hs_code}): {e}")
        return pd.DataFrame()

def calculate_extra_eu(df):
    """
    World 데이터에서 EU27 국가들의 데이터를 뺀 값을 계산하여 반환
    """
    try:
        df['primaryValue'] = pd.to_numeric(df['primaryValue'], errors='coerce').fillna(0)
        df['partnerCode'] = df['partnerCode'].astype(str)
        
        df_world = df[df['partnerCode'] == '0'].copy()
        df_eu = df[df['partnerCode'].isin(EU27_LIST)].copy()
        
        group_cols = ['reporterCode', 'reporterDesc', 'period', 'flowCode', 'flowDesc', 'cmdCode', 'cmdDesc']
        
        # EU 합계
        df_eu_sum = df_eu.groupby(group_cols)['primaryValue'].sum().reset_index()
        df_eu_sum = df_eu_sum.rename(columns={'primaryValue': 'euValue'})
        
        # 병합 및 차감
        merged = pd.merge(df_world, df_eu_sum, on=group_cols, how='left')
        merged['euValue'] = merged['euValue'].fillna(0)
        merged['extraEuValue'] = merged['primaryValue'] - merged['euValue']
        
        merged['primaryValue'] = merged['extraEuValue']
        merged['partnerCode'] = 'EXTRA_EU'
        merged['partnerDesc'] = 'EU27 Extra (Calculated)'
        
        return merged.drop(columns=['euValue', 'extraEuValue'])

    except Exception as e:
        print(f"Calculation Error: {e}")
        return pd.DataFrame()

def preprocess_dataframe(df, original_hs_codes):
    """
    다운로드용 데이터프레임 전처리:
    - 필요한 열만 선택 및 정리
    - 국가명 영문 열 추가 (COUNTRY_NAMES 딕셔너리 활용)
    - cmdCode를 원본 형식 유지 (앞에 0 추가)
    - netWgt, primaryValue 열명에 단위 표시
    """
    if df.empty:
        return df
    
    result = df.copy()
    
    # 국가 코드를 영문 국가명으로 변환하는 함수
    def get_country_name(code):
        code_str = str(code).strip()
        # 먼저 그대로 찾기
        if code_str in COUNTRY_NAMES:
            return COUNTRY_NAMES[code_str]
        # 앞에 0을 붙여서 찾기 (3자리로)
        padded_code = code_str.zfill(3)
        if padded_code in COUNTRY_NAMES:
            return COUNTRY_NAMES[padded_code]
        # 찾지 못하면 코드 그대로 반환
        return code_str
    
    # reporterCode와 partnerCode에서 영문 국가명 생성
    if 'reporterCode' in result.columns:
        result['reporterName'] = result['reporterCode'].apply(get_country_name)
    
    if 'partnerCode' in result.columns:
        result['partnerName'] = result['partnerCode'].apply(get_country_name)
    
    # cmdCode를 원본 HS 코드 형식으로 변환 (앞에 0 추가)
    hs_code_map = {code.lstrip('0'): code for code in original_hs_codes if code}
    hs_code_map.update({code: code for code in original_hs_codes if code})  # 원본도 매핑
    
    def format_cmdcode(code):
        code_str = str(code).strip()
        # 먼저 원본 매핑 확인
        if code_str in hs_code_map:
            return hs_code_map[code_str]
        # 숫자로 변환 후 매핑 확인
        code_stripped = code_str.lstrip('0')
        if code_stripped in hs_code_map:
            return hs_code_map[code_stripped]
        return code_str
    
    result['cmdCode'] = result['cmdCode'].apply(format_cmdcode)
    
    # 필요한 열 선택 및 순서 정렬
    columns_to_keep = [
        'period',
        'reporterCode', 'reporterName',
        'partnerCode', 'partnerName',
        'flowCode',
        'cmdCode',
        'netWgt', 'primaryValue'
    ]
    
    # 존재하는 열만 선택
    available_cols = [col for col in columns_to_keep if col in result.columns]
    result = result[available_cols]
    
    # 열 이름 변경: 단위 표시
    rename_map = {
        'netWgt': 'netWgt (kg)',
        'primaryValue': 'primaryValue (USD)'
    }
    result = result.rename(columns=rename_map)
    
    return result

def create_alluvial_diagram(df, font_size=20,
                            reporter_color='#2E86AB', 
                            hscode_color='#A23B72', partner_color='#F18F01',
                            reporter_font_color='#000000',
                            hscode_font_color='#000000',
                            partner_font_color='#000000',
                            link_opacity=0.3, diagram_height=600, 
                            node_thickness=20,
                            group_by_continent=False, 
                            custom_title="",
                            merge_eu27_reporter=False,
                            show_hscode_percentage=False,
                            show_partner_percentage=False,
                            top_n_partners=None):
    """
    Plotly Sankey diagram 생성
    Reporter → cmdCode → Partner (두께: netWgt)
    
    Parameters:
    - font_size: 폰트 크기 (기본값: 20)
    - reporter_color: Reporter 노드 색상
    - hscode_color: HS Code 노드 색상
    - partner_color: Partner 노드 색상
    - reporter_font_color: Reporter 폰트 색상
    - hscode_font_color: HS Code 폰트 색상
    - partner_font_color: Partner 폰트 색상
    - link_opacity: 링크 투명도 (0~1)
    - diagram_height: 다이어그램 높이 (px)
    - node_thickness: 노드 두께 (기본값: 20)
    - group_by_continent: True면 국가를 대륙별로 그룹화
    - custom_title: 제목 문자열 (빈 문자열이면 표시 안함)
    - merge_eu27_reporter: True면 EU27 국가 Reporter를 "EU27"로 통합
    - show_hscode_percentage: True면 HS Code에 비율 표시
    - show_partner_percentage: True면 Partner에 비율 표시
    - top_n_partners: 상위 N개국만 표시, 나머지는 "기타"로 그룹화 (None이면 전체 표시)
    """
    import plotly.graph_objects as go
    
    if df.empty or 'netWgt (kg)' not in df.columns:
        return None
    
    # 결측치 제거 및 데이터 정리
    df_clean = df.copy()
    df_clean['netWgt (kg)'] = pd.to_numeric(df_clean['netWgt (kg)'], errors='coerce').fillna(0)
    df_clean = df_clean[df_clean['netWgt (kg)'] > 0]
    
    if df_clean.empty:
        return None
    
    # EU27 Reporter 통합 (EU27 국가를 "EU27"로 표시)
    if merge_eu27_reporter and 'reporterCode' in df_clean.columns:
        df_clean['reporterName'] = df_clean['reporterCode'].apply(
            lambda x: "EU27" if str(x).zfill(3) in EU27_LIST else df_clean.loc[df_clean['reporterCode'] == x, 'reporterName'].iloc[0] if len(df_clean.loc[df_clean['reporterCode'] == x, 'reporterName']) > 0 else str(x)
        )
        # 더 간단한 방식으로 재작성
        def get_reporter_display(row):
            code = str(row['reporterCode']).zfill(3)
            if code in EU27_LIST:
                return "EU27"
            return row['reporterName']
        df_clean['reporterName'] = df_clean.apply(get_reporter_display, axis=1)
    
    # 대륙별 그룹화 적용 (유럽을 Intra/Extra-EU27로 분리)
    if group_by_continent and 'partnerCode' in df_clean.columns:
        def get_continent_with_eu_split(code):
            """유럽을 Intra-EU27과 Extra-EU27로 분리"""
            code_str = str(code).zfill(3)
            if code_str in EU27_LIST:
                return "Intra-EU27"
            elif code_str in CONTINENT_EUROPE:
                return "Extra-EU27"
            # 기존 대륙 분류
            if code_str in CONTINENT_AFRICA: return "Africa"
            if code_str in CONTINENT_MIDDLE_EAST: return "Middle East"
            if code_str in CONTINENT_EAST_ASIA: return "East Asia"
            if code_str in CONTINENT_SOUTHEAST_ASIA: return "Southeast Asia"
            if code_str in CONTINENT_NORTH_AMERICA: return "North America"
            if code_str in CONTINENT_CENTRAL_SOUTH_AMERICA: return "Central/South America"
            if code_str in CONTINENT_OCEANIA: return "Oceania"
            return "Others"
        
        df_clean['partnerContinent'] = df_clean['partnerCode'].apply(get_continent_with_eu_split)
        # 대륙별로 물량 합산
        df_grouped = df_clean.groupby(['reporterName', 'cmdCode', 'partnerContinent'])['netWgt (kg)'].sum().reset_index()
        df_grouped = df_grouped.rename(columns={'partnerContinent': 'partnerName'})
        df_clean = df_grouped
    
    # 노드 목록 생성
    reporters = df_clean['reporterName'].unique().tolist()
    cmdcodes = df_clean['cmdCode'].unique().tolist()
    
    # Partner를 물량 기준 내림차순 정렬
    partner_volumes = df_clean.groupby('partnerName')['netWgt (kg)'].sum().sort_values(ascending=False)
    partners = partner_volumes.index.tolist()
    
    # 상위 N개국만 표시, 나머지는 "기타"로 그룹화
    if top_n_partners is not None and len(partners) > top_n_partners:
        top_partners = partners[:top_n_partners]
        other_partners = partners[top_n_partners:]
        
        # 데이터프레임에서 "기타" 그룹으로 변환
        df_clean = df_clean.copy()
        df_clean['partnerName'] = df_clean['partnerName'].apply(
            lambda x: x if x in top_partners else '기타'
        )
        # "기타"를 합산하여 재계산
        df_clean = df_clean.groupby(['reporterName', 'cmdCode', 'partnerName'])['netWgt (kg)'].sum().reset_index()
        
        # Partner 다시 정렬
        partner_volumes = df_clean.groupby('partnerName')['netWgt (kg)'].sum().sort_values(ascending=False)
        partners = partner_volumes.index.tolist()
    
    # 전체 물량 계산 (비율 계산용)
    total_volume = df_clean['netWgt (kg)'].sum()
    
    # HS Code별 물량 계산
    hscode_volumes = df_clean.groupby('cmdCode')['netWgt (kg)'].sum()
    
    # cmdCode에 접두사 추가 (비율 포함 여부에 따라)
    cmdcodes_prefixed = []
    for c in cmdcodes:
        if show_hscode_percentage and total_volume > 0:
            pct = (hscode_volumes.get(c, 0) / total_volume) * 100
            cmdcodes_prefixed.append(f"HS-{c}\n({pct:.1f}%)")
        else:
            cmdcodes_prefixed.append(f"HS-{c}")
    
    # Partner 레이블 (비율 포함 여부에 따라)
    partners_labeled = []
    for p in partners:
        if show_partner_percentage and total_volume > 0:
            pct = (partner_volumes.get(p, 0) / total_volume) * 100
            partners_labeled.append(f"{p}\n({pct:.1f}%)")
        else:
            partners_labeled.append(p)
    
    all_nodes = reporters + cmdcodes_prefixed + partners_labeled
    node_indices = {node: i for i, node in enumerate(all_nodes)}
    
    # 원본 cmdCode와 레이블 매핑
    cmdcode_to_label = {c: cmdcodes_prefixed[i] for i, c in enumerate(cmdcodes)}
    partner_to_label = {p: partners_labeled[i] for i, p in enumerate(partners)}
    
    # 링크 1: Reporter → cmdCode
    link1 = df_clean.groupby(['reporterName', 'cmdCode'])['netWgt (kg)'].sum().reset_index()
    sources1 = [node_indices[r] for r in link1['reporterName']]
    targets1 = [node_indices[cmdcode_to_label[c]] for c in link1['cmdCode']]
    values1 = link1['netWgt (kg)'].tolist()
    
    # 링크 2: cmdCode → Partner
    link2 = df_clean.groupby(['cmdCode', 'partnerName'])['netWgt (kg)'].sum().reset_index()
    sources2 = [node_indices[cmdcode_to_label[c]] for c in link2['cmdCode']]
    targets2 = [node_indices[partner_to_label[p]] for p in link2['partnerName']]
    values2 = link2['netWgt (kg)'].tolist()
    
    # 모든 링크 결합
    sources = sources1 + sources2
    targets = targets1 + targets2
    values = values1 + values2
    
    # 대륙별 색상 매핑 (Intra/Extra-EU27 포함)
    continent_colors = {
        "Intra-EU27": "#1f77b4",      # 파란색 (EU 역내)
        "Extra-EU27": "#17becf",      # 청록색 (EU 역외 유럽)
        "Europe": "#1f77b4",          # 파란색 (기본)
        "Africa": "#ff7f0e",          # 주황색
        "Middle East": "#2ca02c",     # 녹색
        "East Asia": "#d62728",       # 빨간색
        "Southeast Asia": "#9467bd",  # 보라색
        "North America": "#8c564b",   # 갈색
        "Central/South America": "#e377c2",  # 분홍색
        "Oceania": "#7f7f7f",         # 회색
        "Others": "#bcbd22"           # 올리브색
    }
    
    # 노드 색상 및 폰트 색상 설정
    node_colors = []
    font_colors = []
    for node in all_nodes:
        if node in reporters:
            node_colors.append(reporter_color)
            font_colors.append(reporter_font_color)
        elif node.startswith('HS-'):
            node_colors.append(hscode_color)
            font_colors.append(hscode_font_color)
        else:
            # Partner 노드 (대륙 체크)
            matched_continent = False
            for continent_name, continent_color in continent_colors.items():
                if node.startswith(continent_name):
                    node_colors.append(continent_color)
                    font_colors.append(partner_font_color)
                    matched_continent = True
                    break
            if not matched_continent:
                node_colors.append(partner_color)
                font_colors.append(partner_font_color)
    
    # 링크 색상 (회색 + 투명도)
    link_color = f'rgba(100, 100, 100, {link_opacity})'
    
    # Sankey 다이어그램 생성 (Plotly Sankey는 노드별 폰트 색상 미지원)
    fig = go.Figure(data=[go.Sankey(
        textfont=dict(size=font_size, color=reporter_font_color),
        node=dict(
            pad=15,
            thickness=node_thickness,
            line=dict(color="black", width=0.5),
            label=all_nodes,
            color=node_colors
        ),
        link=dict(
            source=sources,
            target=targets,
            value=values,
            color=link_color
        )
    )])
    
    fig.update_layout(
        title_text=custom_title,
        font=dict(size=font_size),
        height=diagram_height
    )
    
    return fig


def render_country_selection_ui(selection_type="Reporter"):
    """
    UN M49 지역별 체크박스 UI 렌더링
    
    Args:
        selection_type: "Reporter" 또는 "Partner"
    
    Returns:
        selected_codes: 선택된 국가 코드 리스트
    """
    selected_codes = []
    
    st.write(f"**{selection_type} 국가 선택** (지역별로 선택)")
    
    # 각 지역별로 expander 생성
    for region_name, sub_regions in COUNTRIES_BY_REGION.items():
        with st.expander(f"🌍 {region_name}"):
            # 지역 전체 선택 체크박스
            region_select_all = st.checkbox(
                f"✓ {region_name} 전체 선택", 
                key=f"{selection_type}_{region_name}_all"
            )
            
            # 하위 지역 처리
            def process_sub_region(sub_region_name, countries, parent_key=""):
                """재귀적으로 하위 지역 처리"""
                codes = []
                
                if isinstance(countries, list):
                    # 리스트인 경우: 국가 코드 목록
                    for code in countries:
                        country_name = COUNTRY_NAMES.get(code, f"Unknown ({code})")
                        if region_select_all or st.checkbox(
                            f"□ {country_name} ({code})",
                            key=f"{selection_type}_{parent_key}_{sub_region_name}_{code}"
                        ):
                            codes.append(code)
                elif isinstance(countries, dict):
                    # 딕셔너리인 경우: 하위 지역이 더 있음
                    sub_region_all = st.checkbox(
                        f"▶ {sub_region_name} 전체 선택",
                        key=f"{selection_type}_{parent_key}_{sub_region_name}_all",
                        value=region_select_all
                    ) or region_select_all
                    
                    st.markdown(f"**{sub_region_name}:**")
                    for nested_sub_name, nested_countries in countries.items():
                        nested_codes = process_sub_region(
                            nested_sub_name, 
                            nested_countries, 
                            parent_key=f"{parent_key}_{sub_region_name}"
                        )
                        if sub_region_all:
                            # 하위 지역 전체 선택 시 모든 코드 추가
                            if isinstance(nested_countries, list):
                                codes.extend(nested_countries)
                            else:
                                codes.extend(nested_codes)
                        else:
                            codes.extend(nested_codes)
                
                return codes
            
            # 각 하위 지역 처리
            for sub_region_name, countries in sub_regions.items():
                region_codes = process_sub_region(sub_region_name, countries, parent_key=region_name)
                if region_select_all:
                    # 전체 선택 시 모든 하위 국가 코드 추가
                    def get_all_codes(data):
                        codes = []
                        if isinstance(data, list):
                            codes.extend(data)
                        elif isinstance(data, dict):
                            for value in data.values():
                                codes.extend(get_all_codes(value))
                        return codes
                    selected_codes.extend(get_all_codes(countries))
                else:
                    selected_codes.extend(region_codes)
    
    # 중복 제거
    selected_codes = list(set(selected_codes))
    
    # 선택된 국가 표시
    if selected_codes:
        st.success(f"✓ {len(selected_codes)}개 국가 선택됨")
        with st.expander("선택된 국가 목록 보기"):
            selected_names = [f"{COUNTRY_NAMES.get(c, f'Unknown ({c})')} ({c})" for c in sorted(selected_codes)]
            st.write(", ".join(selected_names))
    else:
        st.warning("⚠️ 국가를 선택해주세요")
    
    return selected_codes


# --- 웹페이지 UI ---
st.set_page_config(page_title="UN Comtrade 데이터 다운로더", layout="wide")

# [수정] 제목 및 작성자 표시
st.title("📦 UN Comtrade 데이터 다운로더")
st.markdown("""
    <div style='text-align: right; margin-top: -20px; color: #888888;'>
        작성자: Myeong suhwan
    </div>
    <hr>
    """, unsafe_allow_html=True)

st.markdown("보고 국가와 상대국을 직접 선택하거나, **EU 역외 교역**을 자동 계산할 수 있습니다.")

# 사이드바
with st.sidebar:
    st.header("🔑 API 설정")
    api_key = st.text_input("Subscription Key", type="password")
    
    st.write("---")
    st.subheader("⚙️ 교역 구분 (Flow)")
    flow_options = st.multiselect(
        "수집할 항목:",
        ["수입 (Import)", "수출 (Export)"],
        default=["수입 (Import)"]
    )
    
    flow_codes = []
    if "수입 (Import)" in flow_options: flow_codes.append("M")
    if "수출 (Export)" in flow_options: flow_codes.append("X")
    final_flow_code = ",".join(flow_codes)
    
    st.write("---")
    st.subheader("🎨 다이어그램 설정")
    
    # 테마 프리셋
    theme_presets = {
        "기본 (Default)": {"reporter": "#2E86AB", "hscode": "#A23B72", "partner": "#F18F01"},
        "다크 모드 (Dark)": {"reporter": "#00D4FF", "hscode": "#FF6B6B", "partner": "#4ECDC4"},
        "파스텔 (Pastel)": {"reporter": "#A8E6CF", "hscode": "#DDA0DD", "partner": "#FFD93D"},
        "비즈니스 (Business)": {"reporter": "#003366", "hscode": "#666666", "partner": "#CC6600"},
        "네온 (Neon)": {"reporter": "#00FFFF", "hscode": "#FF00FF", "partner": "#FFFF00"},
    }
    
    selected_theme = st.selectbox("테마 선택:", list(theme_presets.keys()))
    theme_colors = theme_presets[selected_theme]
    
    with st.expander("세부 설정"):
        diagram_font_size = st.slider("폰트 크기", min_value=15, max_value=50, value=20)
        diagram_height = st.slider("다이어그램 높이 (px)", min_value=400, max_value=1000, value=600, step=50)
        node_thickness = st.slider("노드 두께", min_value=10, max_value=100, value=20, step=5)
        link_opacity = st.slider("링크 투명도", min_value=0.1, max_value=0.8, value=0.3, step=0.1)
        
        st.caption("노드 색상")
        col_c1, col_c2, col_c3 = st.columns(3)
        with col_c1:
            reporter_color = st.color_picker("Reporter", value=theme_colors["reporter"])
        with col_c2:
            hscode_color = st.color_picker("HS Code", value=theme_colors["hscode"])
        with col_c3:
            partner_color = st.color_picker("Partner", value=theme_colors["partner"])
        
        # 폰트 색상 (전체 레이블에 적용 - Plotly Sankey 제한)
        # 폰트 색상 (전체 레이블에 적용 - Plotly Sankey 제한)
        reporter_font_color = st.color_picker("폰트 색상 (레이블)", value="#000000")
        hscode_font_color = reporter_font_color  # 동일 색상 사용
        partner_font_color = reporter_font_color  # 동일 색상 사용
        
        st.caption("제목 설정")
        show_diagram_title = st.checkbox("제목 표시", value=True)
        if show_diagram_title:
            custom_title = st.text_input("제목 입력", value="Alluvial Diagram: Reporter → HS Code → Partner")
        else:
            custom_title = ""
        
        st.caption("비율 표시 (##.#%)")
        show_hscode_percentage = st.checkbox("HS Code 비율 표시", value=False)
        show_partner_percentage = st.checkbox("Partner 비율 표시", value=False)



# 메인 UI
col1, col2 = st.columns([1, 1])

with col1:
    st.subheader("1. 보고 국가 (Reporter)")
    
    # 기존 그룹 옵션 (빠른 선택용)
    st.write("**빠른 선택 옵션:**")
    reporter_quick_select = st.selectbox(
        "그룹 선택 (선택사항)",
        ["선택 안함"] + [k for k in REPORTER_GROUPS.keys() if k != "직접 입력 (Custom)"],
        key="reporter_quick_select"
    )
    
    if reporter_quick_select != "선택 안함":
        reporter_code = REPORTER_GROUPS[reporter_quick_select]
        display_code = (reporter_code[:30] + '...') if len(reporter_code) > 30 else reporter_code
        st.caption(f"Code: {display_code}")
    else:
        # 체크박스 UI로 개별 국가 선택
        st.write("---")
        selected_reporters = render_country_selection_ui("Reporter")
        # 선택된 국가 코드를 쉼표로 구분된 문자열로 변환
        reporter_code = ",".join(selected_reporters) if selected_reporters else ""

with col2:
    st.subheader("2. 상대국 (Partner)")
    
    # 기존 그룹 옵션 (빠른 선택용)
    st.write("**빠른 선택 옵션:**")
    quick_select = st.selectbox(
        "그룹 선택 (선택사항)",
        ["선택 안함"] + [k for k in PARTNER_GROUPS.keys() if k not in ["직접 입력 (Custom)", "--- 대륙별 선택 ---", "--- 기존 선택 ---"]],
        key="partner_quick_select"
    )
    
    if quick_select != "선택 안함" and quick_select != "--- 대륙별 선택 ---" and quick_select != "--- 기존 선택 ---":
        partner_code_val = PARTNER_GROUPS[quick_select]
        if quick_select.startswith("★"):
            st.success("💡 [자동 계산] World - EU27 = EU 역외 실적 산출")
        else:
            display_code = (partner_code_val[:30] + '...') if len(partner_code_val) > 30 else partner_code_val
            st.caption(f"Code: {display_code}")
    else:
        # 체크박스 UI로 개별 국가 선택
        st.write("---")
        selected_partners = render_country_selection_ui("Partner")
        # 선택된 국가 코드를 쉼표로 구분된 문자열로 변환
        partner_code_val = ",".join(selected_partners) if selected_partners else ""


st.subheader("3. 연도 및 품목")
col3, col4 = st.columns([2, 1])
with col3:
    uploaded_file = st.file_uploader("HS 코드 파일 (CSV/TXT)", type=["csv", "txt"])
with col4:
    selected_years = st.multiselect("연도 선택:", YEAR_OPTIONS, default=["2023"])

if st.button("데이터 수집 시작", type="primary"):
    if not api_key or not uploaded_file or not final_flow_code:
        st.warning("설정 정보를 모두 입력해주세요.")
    elif not reporter_code:
        st.error("⚠️ Reporter 국가를 최소 1개 이상 선택해주세요.")
    elif not partner_code_val:
        st.error("⚠️ Partner 국가를 선택하거나 빠른 선택 옵션을 사용해주세요.")
    else:
        # 파일 읽기
        if uploaded_file.name.endswith('.csv'):
            df_input = pd.read_csv(uploaded_file, dtype=str)
            hs_codes = df_input.iloc[:, 0].dropna().tolist()
        else:
            stringio = uploaded_file.getvalue().decode("utf-8")
            hs_codes = [line.strip() for line in stringio.split('\n') if line.strip()]
        
        # 원본 HS 코드 형식 보존 (중복 제거 전)
        original_hs_codes = [c for c in hs_codes if c]
        hs_codes = list(set(original_hs_codes))
        target_years = sorted(selected_years, reverse=True)
        
        # 보고 국가 분할 (안전 요청)
        if "," in reporter_code:
            reporters_list = [r.strip() for r in reporter_code.split(',') if r.strip()]
        else:
            reporters_list = [reporter_code]

        total_tasks = len(hs_codes) * len(target_years) * len(reporters_list)
        st.write(f"📊 총 작업: {total_tasks}회 요청 예정")
        
        all_data = []
        progress_bar = st.progress(0)
        status_text = st.empty()
        current_task = 0
        
        for code in hs_codes:
            for year in target_years:
                for rep in reporters_list:
                    current_task += 1
                    
                    status_text.text(f"Processing... HS:{code} | Year:{year} | Rep:{rep}")
                    
                    df_result = get_comtrade_data(api_key, code, year, rep, partner_code_val, final_flow_code)
                    
                    if not df_result.empty:
                        all_data.append(df_result)
                    
                    progress_bar.progress(current_task / total_tasks)
                    
                    time.sleep(1.2)
        
        status_text.text("✅ 완료!")
        
        if all_data:
            final_df = pd.concat(all_data, ignore_index=True)
            
            # 데이터 전처리 (열 정리)
            final_df = preprocess_dataframe(final_df, original_hs_codes)
            
            st.success(f"수집 성공! 총 {len(final_df)} 건.")
            
            # 미리보기
            st.dataframe(final_df.head())
            
            # 다운로드
            safe_ptn = "Custom" if quick_select == "선택 안함" else quick_select.split("(")[0].strip()
            csv = final_df.to_csv(index=False).encode('utf-8-sig')
            st.download_button(
                label="📥 결과 다운로드 (CSV)",
                data=csv,
                file_name=f"TradeData_{safe_ptn}_{target_years[0]}.csv",
                mime="text/csv",
            )
            
            # 세션 상태에 데이터 저장 (다이어그램 설정 변경 시 재사용)
            st.session_state['final_df'] = final_df
            st.session_state['partner_mode'] = partner_code_val
            st.session_state['reporter_code'] = reporter_code  # EU27 판단용
        else:
            st.warning("데이터가 없습니다.")

# Alluvial Diagram 섹션 (데이터가 있으면 표시 - 설정 변경 시 즉시 반영)
if 'final_df' in st.session_state and not st.session_state['final_df'].empty:
    st.write("---")
    st.subheader("📊 Alluvial Diagram (Reporter → HS Code → Partner)")
    st.caption("두께: 물량 (kg) 기준")
    
    # 대륙별 그룹화 여부 확인
    group_by_continent = st.session_state.get('partner_mode') == 'ALL_CONTINENTS'
    
    # EU27 Reporter 통합 여부 확인 (EU27 전체 선택 시)
    saved_reporter = st.session_state.get('reporter_code', '')
    merge_eu27 = (saved_reporter == EU27_STR)
    
    # 모든 개별 국가 선택 시 상위 10개국만 표시
    is_all_individual = st.session_state.get('partner_mode') == 'all'
    top_n = 10 if is_all_individual else None
    
    try:
        fig = create_alluvial_diagram(
            st.session_state['final_df'],
            font_size=diagram_font_size,
            reporter_color=reporter_color,
            hscode_color=hscode_color,
            partner_color=partner_color,
            reporter_font_color=reporter_font_color,
            hscode_font_color=hscode_font_color,
            partner_font_color=partner_font_color,
            link_opacity=link_opacity,
            diagram_height=diagram_height,
            node_thickness=node_thickness,
            group_by_continent=group_by_continent,
            custom_title=custom_title,
            merge_eu27_reporter=merge_eu27,
            show_hscode_percentage=show_hscode_percentage,
            show_partner_percentage=show_partner_percentage,
            top_n_partners=top_n
        )
        if fig:
            st.plotly_chart(fig, use_container_width=True, key="main_diagram")
            
            # 다운로드용 다이어그램 (제목 없이)
            fig_download = create_alluvial_diagram(
                st.session_state['final_df'],
                font_size=diagram_font_size,
                reporter_color=reporter_color,
                hscode_color=hscode_color,
                partner_color=partner_color,
                reporter_font_color=reporter_font_color,
                hscode_font_color=hscode_font_color,
                partner_font_color=partner_font_color,
                link_opacity=link_opacity,
                diagram_height=diagram_height,
                node_thickness=node_thickness,
                group_by_continent=group_by_continent,
                custom_title="",  # 다운로드용은 제목 없음
                merge_eu27_reporter=merge_eu27,
                show_hscode_percentage=show_hscode_percentage,
                show_partner_percentage=show_partner_percentage,
                top_n_partners=top_n
            )
            
            # 다운로드 버튼들 (PNG, JPG, PPTX)
            try:
                from pptx import Presentation
                from pptx.util import Inches
                import io
                
                # 이미지 생성 (PNG/JPG 공통)
                img_bytes_png = fig_download.to_image(format="png", width=1200, height=diagram_height, scale=2)
                img_bytes_jpg = fig_download.to_image(format="jpeg", width=1200, height=diagram_height, scale=2)
                
                # PPT 생성
                prs = Presentation()
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                img_stream = io.BytesIO(img_bytes_png)
                slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.5), width=Inches(12.333))
                ppt_stream = io.BytesIO()
                prs.save(ppt_stream)
                ppt_bytes = ppt_stream.getvalue()
                
                # 다운로드 버튼 3개 표시
                col_d1, col_d2, col_d3 = st.columns(3)
                with col_d1:
                    st.download_button("📥 PNG", img_bytes_png, "alluvial_diagram.png", "image/png")
                with col_d2:
                    st.download_button("📥 JPG", img_bytes_jpg, "alluvial_diagram.jpg", "image/jpeg")
                with col_d3:
                    st.download_button("📥 PPTX", ppt_bytes, "alluvial_diagram.pptx",
                                       "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            except Exception as download_error:
                st.caption("다운로드를 위해 kaleido, python-pptx 패키지가 필요합니다.")
        else:
            st.info("다이어그램을 생성할 데이터가 충분하지 않습니다. (물량 데이터 필요)")
    except Exception as e:
        st.error(f"다이어그램 생성 오류: {e}")

